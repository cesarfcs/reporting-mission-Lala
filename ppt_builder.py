# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.chart.data import CategoryChartData

DARK = RGBColor(12, 12, 16)
CARD = RGBColor(20, 20, 26)
TEXT = RGBColor(230, 230, 235)
ACCENT_RED = RGBColor(214, 73, 51)
ACCENT_BLUE = RGBColor(67, 133, 255)
ACCENT_GREEN = RGBColor(58, 186, 120)
ACCENT_PURPLE = RGBColor(154, 99, 210)

def _dark_slide(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK

def _title(slide, text):
    shape = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(11.8), Inches(0.9))
    tf = shape.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = TEXT

def _card(slide, x, y, w, h, title, bullets, accent):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    f = box.fill
    f.solid()
    f.fore_color.rgb = CARD
    box.line.color.rgb = RGBColor(40, 40, 48)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = TEXT
    tf.add_paragraph().text = ""
    for b in bullets:
        if not b:
            continue
        q = tf.add_paragraph()
        q.text = "• " + b
        q.font.size = Pt(14)
        q.font.color.rgb = TEXT
    dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + w - 0.25), Inches(y - 0.1), Inches(0.12), Inches(0.12))
    df = dot.fill
    df.solid()
    df.fore_color.rgb = accent
    dot.line.fill.background()

def build_ppt(out_path, *, client_name, logo_path, report_type, report_cycle, offer_name, offer_contacts_target, channels_list, metrics, date_range_label):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)

    # Slide 1 — Titre
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_slide(s)
    _title(s, f"{client_name} — Rapport de mission")
    sub = s.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(9.0), Inches(0.6))
    sub.text_frame.text = f"{report_type} – {report_cycle} – {date_range_label}"
    sub.text_frame.paragraphs[0].font.size = Pt(16)
    sub.text_frame.paragraphs[0].font.color.rgb = TEXT
    if logo_path:
        try:
            s.shapes.add_picture(logo_path, Inches(10.8), Inches(0.5), height=Inches(0.8))
        except Exception:
            pass

    # Slide 2 — Contexte mission
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_slide(s)
    _title(s, "Contexte mission")
    _card(s, 0.6, 1.4, 3.9, 4.2, "Mission", [
        "Prospection Multicanale" if ("Téléphone" in channels_list and "E-mail" in channels_list) else "Prospection",
        "Téléphone" if "Téléphone" in channels_list else "",
        "E-mail" if "E-mail" in channels_list else "",
        "LinkedIn (optionnel)" if "LinkedIn" in channels_list else "",
    ], accent=ACCENT_BLUE)
    _card(s, 4.7, 1.4, 3.4, 4.2, "Objectifs", [
        f"Contacts : {offer_contacts_target or 'à définir'} par cycle",
        f"Offre : {offer_name}",
    ], accent=ACCENT_RED)
    _card(s, 8.3, 1.4, 4.4, 4.2, "Ciblage (extrait)", [
        "Intitulés de poste : via export",
        "Zone/campagnes : via export",
    ], accent=ACCENT_GREEN)

    # Slide 3 — Contacts adressés
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_slide(s)
    _title(s, "Contacts adressés")
    _card(s, 0.6, 1.4, 6.2, 4.2, "Mission", [
        f"Contacts adressés par téléphone : {metrics.get('contacts_phone', 0)}",
        f"Contacts adressés par e-mail : {metrics.get('contacts_email', 0)}",
    ], accent=ACCENT_RED)
    _card(s, 7.1, 1.4, 5.6, 4.2, "E‑mails & RDV", [
        f"Emailés : {metrics.get('contacts_mailed', 0)} — Ouvertures : {metrics.get('contacts_opened', 0)} — Réponses : {metrics.get('contacts_replied', 0)}",
        f"RDV e‑mail : {metrics.get('rdv_email', 0)}  (conv. {(metrics.get('email_conv_rate', 0.0)*100):.1f} %)",
    ], accent=ACCENT_GREEN)

    # Slide 4 — Performance des appels
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_slide(s)
    _title(s, "Performance des appels")
    cdata = CategoryChartData()
    cdata.categories = ["Contacts appelés", "Connectés", "Pitchés"]
    cdata.add_series("Appels", (
        metrics.get("calls_contacted", 0),
        metrics.get("calls_connected", 0),
        metrics.get("calls_pitched", 0),
    ))
    chart = s.shapes.add_chart(5, Inches(0.7), Inches(1.5), Inches(7.2), Inches(3.6), cdata).chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.show_value = True
    plot.data_labels.number_format = "0"
    _card(s, 8.3, 1.5, 4.4, 3.6, "Analyse", [
        f"{metrics.get('calls_contacted', 0)} contacts appelés (Meeting/No answer/Numéro Faux/Pitch/Sans Suite/Standard)",
        f"{metrics.get('calls_connected', 0)} connectés (Meeting/Pitch/Sans Suite/Standard)",
        f"{metrics.get('calls_pitched', 0)} pitchs — taux {(metrics.get('pitch_rate',0.0)*100):.1f} %",
        f"RDV téléphone : {metrics.get('rdv_phone',0)} — conv. {(metrics.get('calls_conv_rate',0.0)*100):.1f} %",
    ], accent=ACCENT_BLUE)

    # Slide 5 — Performance e‑mailing
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_slide(s)
    _title(s, "Performance e‑mailing")
    cdata = CategoryChartData()
    cdata.categories = ["Emailés", "Ouvertures", "Réponses"]
    cdata.add_series("E‑mail", (
        metrics.get("contacts_mailed", 0),
        metrics.get("contacts_opened", 0),
        metrics.get("contacts_replied", 0),
    ))
    chart = s.shapes.add_chart(5, Inches(0.7), Inches(1.5), Inches(7.2), Inches(3.6), cdata).chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.show_value = True
    plot.data_labels.number_format = "0"
    _card(s, 8.3, 1.5, 4.4, 3.6, "Analyse", [
        f"{metrics.get('contacts_mailed',0)} e‑mailés",
        f"{metrics.get('contacts_opened',0)} ouvertures — {metrics.get('contacts_replied',0)} réponses",
        f"RDV e‑mail : {metrics.get('rdv_email',0)} — conv. {(metrics.get('email_conv_rate',0.0)*100):.1f} %",
    ], accent=ACCENT_GREEN)

    # Slide 6 — RDV obtenus
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_slide(s)
    _title(s, "Rendez‑vous obtenus")
    cdata = CategoryChartData()
    cdata.categories = ["RDV téléphone", "RDV e‑mail"]
    cdata.add_series("RDV", (metrics.get("rdv_phone",0), metrics.get("rdv_email",0)))
    chart = s.shapes.add_chart(5, Inches(1.2), Inches(2.2), Inches(5.5), Inches(2.2), cdata).chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.show_value = True
    plot.data_labels.number_format = "0"
    _card(s, 7.3, 2.0, 5.4, 2.6, "Lecture", [
        f"Téléphone : {metrics.get('rdv_phone',0)} — conv. {(metrics.get('calls_conv_rate',0.0)*100):.1f} %",
        f"E‑mail : {metrics.get('rdv_email',0)} — conv. {(metrics.get('email_conv_rate',0.0)*100):.1f} %",
    ], accent=ACCENT_PURPLE)

    prs.save(out_path)
